from __future__ import annotations

from pathlib import Path

import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
import seaborn as sns
from matplotlib import font_manager
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


BASE_DIR = Path(__file__).resolve().parent
DATA_PATH = BASE_DIR / "data" / "kbo_attendance.csv"
ARTIFACT_DIR = BASE_DIR / "artifacts"
PPT_PATH = BASE_DIR / "KBO_관중예측_LLM스마트티켓팅_발표자료_디자인개선.pptx"

WIDE_W = Inches(13.333)
WIDE_H = Inches(7.5)

COLORS = {
    "ink": RGBColor(18, 31, 44),
    "muted": RGBColor(89, 103, 116),
    "green": RGBColor(31, 111, 88),
    "green_dark": RGBColor(22, 76, 67),
    "clay": RGBColor(191, 122, 79),
    "blue": RGBColor(49, 93, 125),
    "bg": RGBColor(246, 248, 246),
    "paper": RGBColor(251, 252, 249),
    "card": RGBColor(255, 255, 255),
    "line": RGBColor(214, 224, 220),
    "line_dark": RGBColor(174, 194, 187),
    "soft_green": RGBColor(231, 241, 237),
    "soft_clay": RGBColor(246, 235, 226),
    "soft_blue": RGBColor(232, 241, 246),
    "cream": RGBColor(250, 246, 238),
}


def setup_font() -> str:
    available_fonts = {f.name for f in font_manager.fontManager.ttflist}
    for name in ["Malgun Gothic", "NanumGothic", "Noto Sans CJK KR", "Noto Sans KR"]:
        if name in available_fonts:
            plt.rcParams["font.family"] = name
            break
    plt.rcParams["axes.unicode_minus"] = False
    sns.set_theme(style="whitegrid", rc={"font.family": plt.rcParams["font.family"], "axes.unicode_minus": False})
    return plt.rcParams["font.family"][0] if isinstance(plt.rcParams["font.family"], list) else plt.rcParams["font.family"]


def money_fmt(x, _pos=None):
    return f"{int(x):,}"


def load_data() -> pd.DataFrame:
    df = pd.read_csv(DATA_PATH, encoding="utf-8-sig")
    df["date"] = pd.to_datetime(df["date"])
    df["attendance"] = df["attendance"].astype(float)
    real = df.dropna(subset=["attendance"]).copy()
    real["month"] = real["date"].dt.month
    real["weekday_num"] = real["date"].dt.weekday
    real["is_weekend"] = (real["weekday_num"] >= 5).astype(int)
    return real


def make_charts(df: pd.DataFrame) -> dict[str, Path]:
    chart_dir = ARTIFACT_DIR / "ppt_charts"
    chart_dir.mkdir(exist_ok=True)
    chart_paths: dict[str, Path] = {}

    monthly = df.groupby(["season", "month"], as_index=False)["attendance"].mean()
    plt.figure(figsize=(8.7, 4.4))
    sns.lineplot(data=monthly, x="month", y="attendance", hue="season", marker="o", linewidth=2.4)
    plt.title("시즌별 월 평균 관중 수", fontsize=16, weight="bold")
    plt.xlabel("월")
    plt.ylabel("평균 관중 수")
    plt.gca().yaxis.set_major_formatter(money_fmt)
    plt.tight_layout()
    path = chart_dir / "monthly_attendance.png"
    plt.savefig(path, dpi=180, bbox_inches="tight")
    plt.close()
    chart_paths["monthly"] = path

    weekday_order = ["월", "화", "수", "목", "금", "토", "일"]
    weekday_avg = (
        df.groupby("weekday", as_index=False)["attendance"]
        .mean()
        .assign(weekday=lambda x: pd.Categorical(x["weekday"], weekday_order, ordered=True))
        .sort_values("weekday")
    )
    plt.figure(figsize=(7.5, 4.2))
    ax = sns.barplot(data=weekday_avg, x="weekday", y="attendance", palette="viridis")
    plt.title("요일별 평균 관중 수", fontsize=16, weight="bold")
    plt.xlabel("요일")
    plt.ylabel("평균 관중 수")
    ax.yaxis.set_major_formatter(money_fmt)
    plt.tight_layout()
    path = chart_dir / "weekday_attendance.png"
    plt.savefig(path, dpi=180, bbox_inches="tight")
    plt.close()
    chart_paths["weekday"] = path

    home_avg = df.groupby("home_team", as_index=False)["attendance"].mean().sort_values("attendance", ascending=False)
    plt.figure(figsize=(8, 4.2))
    ax = sns.barplot(data=home_avg, x="home_team", y="attendance", palette="crest")
    plt.title("홈팀별 평균 관중 수", fontsize=16, weight="bold")
    plt.xlabel("홈팀")
    plt.ylabel("평균 관중 수")
    ax.yaxis.set_major_formatter(money_fmt)
    plt.tight_layout()
    path = chart_dir / "home_team_attendance.png"
    plt.savefig(path, dpi=180, bbox_inches="tight")
    plt.close()
    chart_paths["home"] = path

    stadium_capacity = {
        "잠실": 23750,
        "대구": 24000,
        "광주": 20500,
        "문학": 23000,
        "창원": 18128,
        "수원": 18700,
        "고척": 16000,
        "대전": 17000,
        "사직": 22758,
        "청주": 12000,
    }
    heat_df = df.copy()
    heat_df["capacity"] = heat_df["stadium"].map(stadium_capacity)
    heat_df["capacity"] = heat_df["capacity"].fillna(heat_df.groupby("stadium")["attendance"].transform("max"))
    heat_df["attendance_rate_pct"] = heat_df["attendance"] / heat_df["capacity"] * 100
    stadium_month = heat_df.pivot_table(index="stadium", columns="month", values="attendance_rate_pct", aggfunc="mean")
    plt.figure(figsize=(8.7, 4.6))
    sns.heatmap(stadium_month, annot=True, fmt=".0f", cmap="YlGnBu", linewidths=0.4)
    plt.title("구장별/월별 평균 좌석 점유율(%)", fontsize=16, weight="bold")
    plt.xlabel("월")
    plt.ylabel("구장")
    plt.tight_layout()
    path = chart_dir / "stadium_month_heatmap.png"
    plt.savefig(path, dpi=180, bbox_inches="tight")
    plt.close()
    chart_paths["heatmap"] = path

    compare_path = ARTIFACT_DIR / "model_compare.csv"
    compare = pd.read_csv(compare_path) if compare_path.exists() else pd.DataFrame()
    if not compare.empty:
        metric = compare[["model", "mae", "rmse"]].melt("model", var_name="metric", value_name="score")
        plt.figure(figsize=(7.8, 4.3))
        ax = sns.barplot(data=metric, x="model", y="score", hue="metric", palette=["#1f6f58", "#bf7a4f"])
        plt.title("딥러닝 모델 오차 비교", fontsize=16, weight="bold")
        plt.xlabel("모델")
        plt.ylabel("오차")
        ax.yaxis.set_major_formatter(money_fmt)
        plt.tight_layout()
        path = chart_dir / "model_compare.png"
        plt.savefig(path, dpi=180, bbox_inches="tight")
        plt.close()
        chart_paths["models"] = path

    candidate_path = ARTIFACT_DIR / "llm_candidate_games.csv"
    if candidate_path.exists():
        cand = pd.read_csv(candidate_path)
        cand["date"] = pd.to_datetime(cand["date"])
        cand["match"] = cand.apply(lambda r: f"{r['date'].date()} {r['home_team']} vs {r['away_team']}", axis=1)
        cand = cand.sort_values("attendance", ascending=True).tail(8)
        plt.figure(figsize=(8.7, 4.4))
        ax = sns.barplot(data=cand, y="match", x="attendance", hue="crowd_level", dodge=False, palette="Spectral")
        plt.title("LLM 검색 기반 추천 후보 예시", fontsize=16, weight="bold")
        plt.xlabel("관중 수")
        plt.ylabel("후보 경기")
        ax.xaxis.set_major_formatter(money_fmt)
        plt.legend(title="혼잡도", loc="lower right")
        plt.tight_layout()
        path = chart_dir / "llm_candidates.png"
        plt.savefig(path, dpi=180, bbox_inches="tight")
        plt.close()
        chart_paths["candidates"] = path

    return chart_paths


def add_textbox(slide, x, y, w, h, text, size=20, bold=False, color="ink", align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "Malgun Gothic"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = COLORS[color] if isinstance(color, str) else color
    return box


def add_bullets(slide, x, y, w, h, items, size=18, color="ink"):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = "Malgun Gothic"
        p.font.size = Pt(size)
        p.font.color.rgb = COLORS[color]
        p.space_after = Pt(9)
    return box


def add_rect(slide, x, y, w, h, fill="card", line="line", rounded=True):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS[fill]
    shape.line.color.rgb = COLORS[line]
    return shape


def add_header(slide, title, page, total):
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = COLORS["paper"]

    rail = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.18), WIDE_H)
    rail.fill.solid()
    rail.fill.fore_color.rgb = COLORS["green_dark"]
    rail.line.color.rgb = COLORS["green_dark"]

    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, WIDE_W, Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLORS["green"]
    bar.line.color.rgb = COLORS["green"]

    add_rect(slide, 0.56, 0.34, 3.65, 0.34, "soft_green", "line", True)
    add_textbox(slide, 0.72, 0.42, 3.3, 0.16, "KBO 관중 예측 기반 LLM 스마트 티켓팅", 9, True, "green_dark", PP_ALIGN.CENTER)
    add_textbox(slide, 0.58, 0.78, 10.7, 0.55, title, 25, True, "ink")

    add_rect(slide, 12.05, 6.83, 0.48, 0.34, "green_dark", "green_dark", True)
    add_textbox(slide, 12.05, 6.91, 0.48, 0.14, str(page), 9, True, RGBColor(255, 255, 255), PP_ALIGN.CENTER)

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.58), Inches(1.43), Inches(12.15), Inches(0.012))
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS["line"]
    line.line.color.rgb = COLORS["line"]


def add_image(slide, path: Path, x, y, w, h):
    slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))


def cover_slide(prs, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = COLORS["paper"]

    rail = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.2), WIDE_H)
    rail.fill.solid()
    rail.fill.fore_color.rgb = COLORS["green_dark"]
    rail.line.color.rgb = COLORS["green_dark"]

    add_rect(slide, 0.68, 0.68, 12.0, 5.98, "card", "line_dark")
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.68), Inches(0.68), Inches(12.0), Inches(0.2))
    band.fill.solid()
    band.fill.fore_color.rgb = COLORS["green"]
    band.line.color.rgb = COLORS["green"]

    add_rect(slide, 1.16, 1.16, 2.1, 0.38, "soft_green", "line", True)
    add_textbox(slide, 1.31, 1.25, 1.8, 0.16, "AI 서비스 프로젝트", 9, True, "green_dark", PP_ALIGN.CENTER)
    add_textbox(slide, 1.16, 1.78, 10.6, 1.15, "KBO 관중 수 예측 기반\nLLM 스마트 티켓팅 가이드", 36, True, "ink")
    add_textbox(slide, 1.2, 3.18, 9.9, 0.42, "딥러닝 예측 결과를 경기 추천과 예매 의사결정으로 확장한 AI 서비스", 17, False, "muted")

    diamond = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, Inches(10.5), Inches(1.45), Inches(1.55), Inches(1.55))
    diamond.fill.solid()
    diamond.fill.fore_color.rgb = COLORS["cream"]
    diamond.line.color.rgb = COLORS["clay"]
    add_textbox(slide, 10.72, 2.05, 1.12, 0.22, "KBO", 16, True, "clay", PP_ALIGN.CENTER)

    labels = [("Dense 예측", "soft_green", "green_dark"), ("FAISS 검색", "soft_clay", "clay"), ("Ollama LLM", "soft_blue", "blue"), ("티켓팅 추천", "soft_green", "green_dark")]
    for i, (text, fill, color) in enumerate(labels):
        add_rect(slide, 1.2 + i * 2.45, 4.05, 2.05, 0.58, fill, "line")
        add_textbox(slide, 1.2 + i * 2.45, 4.2, 2.05, 0.22, text, 14, True, color, PP_ALIGN.CENTER)
    add_textbox(slide, 1.2, 5.55, 4.7, 0.35, "2026.04.20.", 15, True, "green")
    add_textbox(slide, 1.2, 5.96, 4.7, 0.28, "권우영", 13, False, "muted")
    add_rect(slide, 12.05, 6.83, 0.48, 0.34, "green_dark", "green_dark", True)
    add_textbox(slide, 12.05, 6.91, 0.48, 0.14, "1", 9, True, RGBColor(255, 255, 255), PP_ALIGN.CENTER)


def build_ppt():
    setup_font()
    df = load_data()
    charts = make_charts(df)
    compare = pd.read_csv(ARTIFACT_DIR / "model_compare.csv")
    best = compare.sort_values("mae").iloc[0]

    prs = Presentation()
    prs.slide_width = WIDE_W
    prs.slide_height = WIDE_H
    total = 14

    cover_slide(prs, total)

    slides = [
        ("목차", "agenda"),
        ("프로젝트 배경과 목표", "background"),
        ("데이터 구성과 전처리", "data"),
        ("딥러닝 예측 모델 구성", "model"),
        ("모델 성능 비교", "performance"),
        ("관중 패턴 시각화 인사이트", "insight"),
        ("LLM 확장 전략", "llm_strategy"),
        ("RAG 기반 추천 구조", "rag"),
        ("LLM 추천 후보 예시", "candidate"),
        ("서비스 활용 시나리오", "scenario"),
        ("프로토타입 시연 화면", "prototype"),
        ("개선 및 한계점", "limits"),
        ("결론", "conclusion"),
    ]

    for idx, (title, kind) in enumerate(slides, start=2):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_header(slide, title, idx, total)

        if kind == "agenda":
            add_bullets(slide, 0.95, 1.85, 5.9, 4.4, [
                "1. 프로젝트 배경과 목표",
                "2. 데이터 구성과 전처리",
                "3. 딥러닝 관중 수 예측 모델",
                "4. 시각화 기반 관중 패턴 분석",
                "5. LLM/RAG 기반 스마트 티켓팅 확장",
                "6. 프로토타입, 한계점, 결론",
            ], 21)
            add_rect(slide, 7.25, 1.95, 4.55, 3.35, "card", "line")
            add_textbox(slide, 7.65, 2.35, 3.75, 0.45, "발표 핵심", 20, True, "green", PP_ALIGN.CENTER)
            add_textbox(slide, 7.65, 3.05, 3.75, 1.35, "관중 수를 예측하는 데서 끝나지 않고,\n사용자가 언제 어떤 경기를 예매할지 판단하도록 돕는 AI 서비스로 확장", 19, True, "ink", PP_ALIGN.CENTER)

        elif kind == "background":
            add_bullets(slide, 0.9, 1.8, 5.9, 3.7, [
                "KBO 경기 관중 수는 팀, 구장, 요일, 공휴일, 라이벌전 여부에 따라 크게 달라짐",
                "사용자는 인기 경기와 혼잡도를 예매 전에 미리 알고 싶어 함",
                "기존 딥러닝 예측 결과를 실제 의사결정 서비스로 연결하는 것이 프로젝트의 확장 목표",
            ], 20)
            add_rect(slide, 7.15, 1.95, 4.85, 2.95, "soft_green", "line")
            add_textbox(slide, 7.55, 2.35, 4.05, 0.45, "문제 정의", 22, True, "green", PP_ALIGN.CENTER)
            add_textbox(slide, 7.55, 3.05, 4.05, 1.0, "“이 경기는 얼마나 붐빌까?”\n“언제 예매해야 할까?”", 24, True, "ink", PP_ALIGN.CENTER)

        elif kind == "data":
            add_bullets(slide, 0.85, 1.82, 5.7, 3.8, [
                "2024~2026 KBO 경기 일정 및 실제 관중 데이터 사용",
                "주요 변수: 날짜, 요일, 홈팀, 원정팀, 구장, 관중 수, 공휴일 여부, 시즌",
                "관중 수가 있는 실제 경기만 학습/시각화에 사용",
                "향후 일정은 예측과 추천 서비스의 입력 데이터로 활용",
            ], 19)
            add_image(slide, charts["home"], 6.95, 1.8, 5.3, 3.4)
            add_textbox(slide, 7.05, 5.35, 5.0, 0.55, "노트북 시각화 자료 활용: 홈팀별 평균 관중 수", 13, False, "muted", PP_ALIGN.CENTER)

        elif kind == "model":
            boxes = [("입력 특성", "홈팀, 원정팀, 구장\n월, 요일, 주말\n공휴일, 라이벌전, 시즌"), ("예측 모델", "Dense\nLSTM\nGRU"), ("출력", "예상 관중 수\n좌석 점유율\n예매 혼잡도")]
            for i, (head, body) in enumerate(boxes):
                x = 0.85 + i * 4.15
                add_rect(slide, x, 1.95, 3.55, 3.05, "card", "line")
                add_textbox(slide, x + 0.25, 2.3, 3.05, 0.35, head, 20, True, "green", PP_ALIGN.CENTER)
                add_textbox(slide, x + 0.35, 3.0, 2.85, 1.35, body, 19, False, "ink", PP_ALIGN.CENTER)
            add_textbox(slide, 1.0, 5.45, 11.2, 0.45, "딥러닝 모델은 수치 예측을 담당하고, LLM은 예측 결과의 설명과 추천을 담당하도록 역할을 분리", 18, True, "ink", PP_ALIGN.CENTER)

        elif kind == "performance":
            add_image(slide, charts["models"], 0.8, 1.75, 6.1, 3.75)
            add_rect(slide, 7.35, 1.9, 4.6, 3.1, "card", "line")
            add_textbox(slide, 7.75, 2.25, 3.8, 0.4, "선정 기준", 21, True, "green", PP_ALIGN.CENTER)
            add_textbox(slide, 7.75, 2.9, 3.8, 1.55, f"가장 낮은 MAE를 보인 {best['model']} 모델을 기본 예측 모델로 사용\nMAE 약 {int(round(best['mae'])):,}명", 20, True, "ink", PP_ALIGN.CENTER)

        elif kind == "insight":
            add_image(slide, charts["monthly"], 0.75, 1.75, 5.85, 3.2)
            add_image(slide, charts["weekday"], 6.9, 1.75, 5.4, 3.2)
            add_bullets(slide, 1.0, 5.35, 11.1, 0.85, [
                "월별, 요일별 관중 흐름은 예매 추천의 핵심 근거가 됨",
                "주말/공휴일/인기 팀 경기는 혼잡도와 매진 가능성을 함께 안내해야 함",
            ], 15, "muted")

        elif kind == "llm_strategy":
            steps = [("1", "예측", "Dense 모델로\n예상 관중 산출"), ("2", "검색", "유사 경기와\n과거 관중 패턴 검색"), ("3", "LLM", "질문 의도를 반영해\n추천 답변 생성"), ("4", "서비스", "예매 시점과\n혼잡도 안내")]
            for i, (num, head, body) in enumerate(steps):
                x = 0.72 + i * 3.12
                add_rect(slide, x, 2.0, 2.45, 1.75, "soft_green" if i % 2 == 0 else "soft_clay", "line")
                add_textbox(slide, x + 0.18, 2.2, 0.45, 0.35, num, 16, True, "green", PP_ALIGN.CENTER)
                add_textbox(slide, x + 0.65, 2.18, 1.45, 0.35, head, 18, True, "ink", PP_ALIGN.CENTER)
                add_textbox(slide, x + 0.25, 2.8, 1.95, 0.65, body, 15, False, "ink", PP_ALIGN.CENTER)
            add_textbox(slide, 1.0, 4.65, 11.2, 0.8, "노트북 시각화 자료는 LLM 답변의 근거를 설명하는 보조 자료로 활용하고,\n발표의 중심은 실제 스마트 티켓팅 서비스 흐름에 둠", 20, True, "ink", PP_ALIGN.CENTER)

        elif kind == "rag":
            add_bullets(slide, 0.85, 1.85, 5.6, 3.7, [
                "각 경기 정보를 한 문장짜리 검색 문서로 변환",
                "TF-IDF로 질문과 문서의 텍스트 유사도를 계산",
                "FAISS로 유사 경기 후보를 빠르게 검색",
                "검색 결과와 예측값을 LLM 프롬프트에 넣어 추천 답변 생성",
            ], 19)
            add_rect(slide, 7.05, 1.95, 4.95, 3.15, "soft_blue", "line")
            add_textbox(slide, 7.45, 2.3, 4.15, 0.4, "예시 질문", 20, True, "blue", PP_ALIGN.CENTER)
            add_textbox(slide, 7.45, 3.0, 4.15, 1.25, "“주말에 가족이 보기 좋고\n관중이 많은 인기 경기 추천해줘”", 21, True, "ink", PP_ALIGN.CENTER)

        elif kind == "prototype":
            add_rect(slide, 0.8, 1.75, 5.55, 4.55, "card", "line")
            add_rect(slide, 6.95, 1.75, 5.55, 4.55, "card", "line")
            add_textbox(slide, 1.05, 2.05, 5.05, 0.5, "캡쳐화면 필요: 예측/티켓팅 가이드", 20, True, "green", PP_ALIGN.CENTER)
            add_textbox(slide, 7.2, 2.05, 5.05, 0.5, "캡쳐화면 필요: AI 상담 결과", 20, True, "green", PP_ALIGN.CENTER)
            add_textbox(slide, 1.25, 3.05, 4.65, 1.6, "Streamlit 화면에서 날짜 선택,\n경기별 예상 관중,\n좌석 점유율과 예매 안내가 보이도록 캡쳐", 18, False, "ink", PP_ALIGN.CENTER)
            add_textbox(slide, 7.4, 3.05, 4.65, 1.6, "예시 질문 입력창,\nLLM 답변,\n추천 근거 경기 표가 함께 보이도록 캡쳐", 18, False, "ink", PP_ALIGN.CENTER)

        elif kind == "candidate":
            add_image(slide, charts["candidates"], 0.82, 1.7, 6.25, 3.85)
            add_bullets(slide, 7.45, 1.9, 4.55, 3.35, [
                "질문과 가까운 경기 후보를 검색해 추천 근거로 사용",
                "후보는 관중 수, 좌석 점유율, 혼잡도 기준으로 해석",
                "LLM 답변은 데이터 기반 근거를 함께 제시하도록 설계",
            ], 19)

        elif kind == "scenario":
            add_bullets(slide, 0.95, 1.85, 5.7, 3.8, [
                "가족 관람: 주말 경기 중 인기와 혼잡도를 함께 비교",
                "팬 관람: 특정 팀 경기 중 매진 가능성이 높은 경기 우선 안내",
                "여유 관람: 평일 또는 점유율 낮은 경기 추천",
                "구단 운영: 예상 관중 기반 매표/운영 인력 계획 보조",
            ], 20)
            add_rect(slide, 7.25, 1.95, 4.55, 3.25, "soft_clay", "line")
            add_textbox(slide, 7.65, 2.35, 3.75, 0.4, "서비스 가치", 21, True, "clay", PP_ALIGN.CENTER)
            add_textbox(slide, 7.65, 3.05, 3.75, 1.35, "예측 모델의 숫자를\n실제 예매 행동으로 연결", 23, True, "ink", PP_ALIGN.CENTER)

        elif kind == "limits":
            add_bullets(slide, 0.9, 1.8, 5.85, 3.9, [
                "날씨, 선발 투수, 팀 순위, 이벤트 정보 같은 실시간 변수가 부족함",
                "LLM 답변은 검색된 데이터에 의존하므로 최신 정보 반영이 필요함",
                "현재는 로컬 Ollama 기반이라 실행 환경에 따라 응답 속도가 달라질 수 있음",
            ], 19)
            add_bullets(slide, 7.05, 1.8, 5.0, 3.9, [
                "개선: 날씨 API와 실시간 잔여 좌석 데이터 연동",
                "개선: 실제 예매 로그 기반 개인화 추천",
                "개선: 경기별 이벤트/상대 전적/순위 정보를 RAG 문서에 추가",
            ], 19, "green_dark")

        elif kind == "conclusion":
            add_rect(slide, 1.15, 1.85, 11.0, 3.95, "card", "line")
            add_textbox(slide, 1.6, 2.35, 10.1, 0.6, "결론", 30, True, "green", PP_ALIGN.CENTER)
            add_textbox(slide, 1.8, 3.25, 9.7, 1.45, "KBO 관중 예측 모델을 기반으로,\n사용자가 이해할 수 있는 추천과 예매 판단을 제공하는\nLLM 스마트 티켓팅 서비스로 확장했다.", 25, True, "ink", PP_ALIGN.CENTER)
            add_textbox(slide, 1.8, 5.2, 9.7, 0.35, "딥러닝은 예측을, LLM은 설명과 추천을 담당한다.", 17, True, "muted", PP_ALIGN.CENTER)

    prs.save(PPT_PATH)
    print(PPT_PATH)


if __name__ == "__main__":
    build_ppt()
