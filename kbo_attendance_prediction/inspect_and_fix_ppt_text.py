from pathlib import Path
from zipfile import ZipFile

from pptx import Presentation
from pptx.util import Pt


BASE_DIR = Path(__file__).resolve().parent
SRC = BASE_DIR / "KBO_관중예측_LLM경기추천가이드_발표자료.pptx"
OUT = BASE_DIR / "KBO_관중예측_LLM경기추천가이드_발표자료_최종수정.pptx"

REPLACEMENTS = {
    "AI 스마트티켓팅가이드": "AI 경기 추천 가이드",
    "AI 스마트 티켓팅 가이드": "AI 경기 추천 가이드",
    "AI 스마트티켓팅 가이드": "AI 경기 추천 가이드",
    "AI 티켓팅상담": "AI 경기 추천 상담",
    "AI 티켓팅 상담": "AI 경기 추천 상담",
    "AI티켓팅상담": "AI 경기 추천 상담",
    "AI티켓팅 상담": "AI 경기 추천 상담",
    "스마트티켓팅가이드": "경기 추천 가이드",
    "스마트티켓팅 가이드": "경기 추천 가이드",
    "스마트 티켓팅 가이드": "경기 추천 가이드",
    "스마트티켓팅": "경기 추천",
    "스마트 티켓팅": "경기 추천",
    "티켓팅상담": "경기 추천 상담",
    "티켓팅 상담": "경기 추천 상담",
    "티켓팅": "경기 추천",
}


def replace_in_shape(shape) -> None:
    if not hasattr(shape, "text_frame"):
        return
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            text = run.text
            for old, new in REPLACEMENTS.items():
                text = text.replace(old, new)
            run.text = text


def fix_pages(prs: Presentation) -> None:
    total = len(prs.slides)
    for slide_idx, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if not hasattr(shape, "text_frame"):
                continue
            text = getattr(shape, "text", "").strip()
            page_tokens = {
                f"{slide_idx:02d}/14",
                f"{slide_idx:02d} / 14",
                f"{slide_idx:02d}/{total:02d}",
                f"{slide_idx:02d} / {total:02d}",
            }
            if text in page_tokens:
                shape.text_frame.clear()
                p = shape.text_frame.paragraphs[0]
                p.alignment = 2
                run = p.add_run()
                run.text = str(slide_idx)
                run.font.name = "Malgun Gothic"
                run.font.size = Pt(9)


def remaining_hits(path: Path) -> list[tuple[str, str]]:
    keys = ["스마트티켓팅", "스마트 티켓팅", "티켓팅", "AI 스마트", "AI 티켓"]
    hits = []
    with ZipFile(path) as z:
        for name in z.namelist():
            if name.startswith("ppt/slides/slide") and name.endswith(".xml"):
                text = z.read(name).decode("utf-8", errors="ignore")
                for key in keys:
                    if key in text:
                        hits.append((name, key))
    return hits


def main() -> None:
    prs = Presentation(SRC)
    for slide in prs.slides:
        for shape in slide.shapes:
            replace_in_shape(shape)
    fix_pages(prs)
    prs.save(OUT)
    print(f"output={OUT}")
    print(f"remaining_hits={remaining_hits(OUT)}")


if __name__ == "__main__":
    main()
