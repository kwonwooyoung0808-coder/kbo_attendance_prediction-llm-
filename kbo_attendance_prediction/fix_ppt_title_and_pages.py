from pathlib import Path

from pptx import Presentation
from pptx.util import Pt


BASE_DIR = Path(__file__).resolve().parent
OUT = BASE_DIR / "KBO_관중예측_LLM경기추천가이드_발표자료.pptx"


def find_original_ppt() -> Path:
    candidates = [
        p
        for p in BASE_DIR.glob("*.pptx")
        if "디자인개선" not in p.name and "경기추천가이드" not in p.name and p.stat().st_size > 100_000
    ]
    if not candidates:
        raise FileNotFoundError("원본 발표자료 PPT를 찾지 못했습니다.")
    return max(candidates, key=lambda p: p.stat().st_size)


def replace_text_preserving_runs(prs: Presentation) -> None:
    replacements = {
        "KBO 관중 예측 기반 LLM 스마트 티켓팅": "KBO 관중 예측 기반 LLM 경기 추천 가이드",
        "LLM 스마트 티켓팅 가이드": "LLM 경기 추천 가이드",
        "스마트 티켓팅": "경기 추천 가이드",
        "티켓팅 추천": "경기 추천",
        "티켓팅": "경기 추천",
    }

    for slide_idx, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if not hasattr(shape, "text_frame"):
                continue

            text = getattr(shape, "text", "").strip()
            if text in {f"{slide_idx:02d}/14", f"{slide_idx:02d} / 14", f"{slide_idx:02d}/{len(prs.slides):02d}", f"{slide_idx:02d} / {len(prs.slides):02d}"}:
                shape.text_frame.clear()
                p = shape.text_frame.paragraphs[0]
                p.alignment = 2
                run = p.add_run()
                run.text = str(slide_idx)
                run.font.name = "Malgun Gothic"
                run.font.size = Pt(9)
                continue

            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    new_text = run.text
                    for old, new in replacements.items():
                        new_text = new_text.replace(old, new)
                    run.text = new_text


def main() -> None:
    src = find_original_ppt()
    prs = Presentation(src)
    replace_text_preserving_runs(prs)
    prs.save(OUT)
    print(f"source={src}")
    print(f"output={OUT}")


if __name__ == "__main__":
    main()
